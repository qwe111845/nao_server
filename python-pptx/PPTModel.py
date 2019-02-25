# -*- coding: utf-8 -*-
from pptx import Presentation
import win32com
import win32com.client
from PIL import Image, ImageTk
import sys
import PPT_Client
import threading as td

reload(sys)
sys.setdefaultencoding('utf8')


class PPTModel:

    def __init__(self):
        self.ppt_file = ''
        self.save_dir = ''
        self.fname = ''
        self.link = ''
        self.ppt_number = 0
        self.ppt_len = 0
        self.normal_mode = True
        self.ppt_ready = False
        self.ppt_text = []
        self.client = PPT_Client.PPTClient()

    def set_img_data(self, save_dir):
        self.save_dir = save_dir
        fstart = self.save_dir.rindex('/')
        self.fname = self.save_dir[fstart + 1:]

        self.save_dir = self.save_dir.replace('/', r'\\')
        import glob

        target = self.save_dir + "\\*.jpg"

        self.ppt_len = len(glob.glob(target))
        self.normal_mode = False

    def get_slide_text(self, text_list):
        slide_texts = ''
        for i, text in enumerate(text_list):
            slide_texts += str(i + 1) + '.\n' + str(text) + '\n'

        return slide_texts

    def set_data(self, ppt_file, save_dir):
        self.normal_mode = True
        self.ppt_file = ppt_file
        self.save_dir = save_dir
        self.ppt_text = []

        fstart = self.ppt_file.rindex('/')
        self.fname = self.ppt_file[fstart + 1:-5]

        self.ppt_file = self.ppt_file.replace('/', r'\\')
        self.save_dir = self.save_dir.replace('/', r'\\')

        self.ppt2png(self.fname, self.ppt_file, self.save_dir)

        self.ppt_ready = True

    def img_set(self, w_box, h_box):
        if self.normal_mode:
            img_location = self.save_dir + r"\\" + self.fname + r"\\投影片" + str(self.ppt_number + 1) + r".jpg"
            text = self.get_slide_text(self.ppt_text[self.ppt_number])
        else:
            img_location = str(self.save_dir) + r"\\" + r"\\投影片" + str(self.ppt_number + 1) + r".jpg"
            # self.text_set(self.ppt_text[self.ppt_number])
        pil_image = Image.open(img_location)
        pil_image_resized = self.resize(w_box, h_box, pil_image)

        tk_image = ImageTk.PhotoImage(pil_image_resized)
        return text, tk_image

    def ppt_start(self):
        self.ppt_number = 0

    def previous_ppt(self):
        self.ppt_number -= 1

    def next_ppt(self):
        self.ppt_number += 1

    def ppt2png(self, ppt_file_name, ppt_path, save_path):

        powerpoint = win32com.client.Dispatch('PowerPoint.Application')
        powerpoint.Visible = True
        ppt = powerpoint.Presentations.Open(ppt_path)
        ppt.SaveAs(save_path + r'\\' + ppt_file_name.rsplit('.')[0], 17)
        ppt.Close()
        powerpoint.Quit()

        prs = Presentation(self.ppt_file)
        self.ppt_len = len(prs.slides)

        for slide in prs.slides:
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = ''
                    for run in paragraph.runs:
                        text += run.text
                    text_runs.append(text)
            self.ppt_text.append(text_runs)

    def resize(self, w_box, h_box, pil_image):
        w, h = pil_image.size
        f1 = 1.0 * w_box / w
        f2 = 1.0 * h_box / h
        factor = min([f1, f2])
        width = int(w * factor)
        height = int(h * factor)
        return pil_image.resize((width, height), Image.ANTIALIAS)

    def folder_setting_completed(self):
        if len(self.save_dir) > 0:
            return True
        else:
            return False

    def ppt_is_ready(self):
        return self.ppt_ready

    def get_ppt_number(self):
        return self.ppt_number

    def get_ppt_len(self):
        return self.ppt_len

    def link_nao(self, ip):
        self.client.main(ip)
        self.link = self.client.get_status()
        return self.link

    def link_success(self):
        self.link = self.client.get_status()
        print(self.link)
        if self.link == '關閉連線' or self.link == '連線失敗' or self.link == "尚未連線":
            return False
        else:
            return True

    def get_link_status(self):
        return self.link

    def nao_speak(self):
        texts = self.ppt_text[self.ppt_number]
        send_text = ''
        for text in texts:
            send_text += text + ":::"

        self.client.set_command(send_text)
        self.link = self.client.get_status()
        return self.link

    def close_connection(self):
        self.client.set_command('Quit ppt mode')
        self.link = self.client.get_status()
        return self.link

    def nao_command(self, command):
        if command == 'Chinese':
            self.client.set_command('Chinese')
        elif command == 'English':
            self.client.set_command('English')
        elif command == 'Increase':
            self.client.set_command('Increased speech rate')
        elif command == 'Reduce':
            self.client.set_command('Reduce speech rate')
        else:
            self.client.set_command(command)


