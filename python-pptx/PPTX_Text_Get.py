from pptx import Presentation
import os
prs = Presentation("robot.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []
import win32com
import win32com.client


def ppt2png(pptFileName, ppt_path, save_path):

    powerpoint = win32com.client.Dispatch('PowerPoint.Application')
    powerpoint.Visible = True
    ppt = powerpoint.Presentations.Open(ppt_path)
    ppt.SaveAs(save_path + pptFileName.rsplit('.')[0], 17)
    ppt.Close()
    powerpoint.Quit()

for slide in prs.slides:
    for shape in slide.shapes:
        text = ''
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text
        text_runs.append(text)
ppt2png('robotdd.pptx',r'C:\Users\lin\PycharmProjects\nao_server\python-pptx\robot.pptx',r'C:\\Users\\lin\\PycharmProjects\\nao_server\\python-pptx\\')

for i, text in enumerate(text_runs):
    print(str(i) + text)