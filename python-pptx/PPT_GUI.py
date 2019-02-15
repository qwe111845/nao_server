# -*- coding: utf-8 -*-

import Tkinter as tk
import tkMessageBox as tm
import os
import tkFileDialog
from pptx import Presentation
import win32com
import win32com.client
from PIL import Image, ImageTk
import socket
import sys
import time

reload(sys)
sys.setdefaultencoding('utf8')


class PPTGUI(tk.Frame):

    def __init__(self, root=None):
        tk.Frame.__init__(self, root)
        self.root = root
        self.var = tk.StringVar()
        self.ppt_file = ''
        self.save_dir = ''
        self.jpg_dir = ''
        self.fname = ''
        self.sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
        self.link = tk.StringVar()
        self.ppt_number = 0
        self.ppt_len = 0
        self.normal_mode = True
        self.ppt_text = []
        self.ip_text = tk.Text(self)
        port = 5007
        s = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
        s.connect(('140.134.26.200', port))
        s.send('get robot port(G)')
        self.robot_port = s.recv(1024)
        time.sleep(1)
        s.close()
        self.init_gui()

    def init_gui(self):
        self.master.title("PPT播放工具")
        self.link.set('尚未連線')
        self.pack(fill='both', expand=1)

        self.ip_text.place(relx=0.85, relheight=0.05, relwidth=0.15)

        link_button = tk.Button(self, text='Link robot', command=self.link_nao)
        link_button.place(relx=0.85, rely=0.05, relheight=0.03, relwidth=0.15)

        self.link_label = tk.Label(self, textvariable=self.link, bg='LightCyan')
        self.link_label.place(relx=0.85, rely=0.08, relheight=0.05, relwidth=0.15)

        ppt_selection_button = tk.Button(self, text="轉換PPT", command=self.set_data)
        ppt_selection_button.place(relx=0.85, rely=0.13, relheight=0.05, relwidth=0.075)

        ppt_selection1_button = tk.Button(self, text="選擇PPT資料夾", command=self.set_img_data)
        ppt_selection1_button.place(relx=0.925, rely=0.13, relheight=0.05, relwidth=0.075)

        start_button = tk.Button(self, text="Start", command=self.ppt_start)
        start_button.place(relx=0.85, rely=0.18, relheight=0.04, relwidth=0.15)

        previous_button = tk.Button(self, text="Previous slide", command=self.previous_ppt)
        previous_button.place(relx=0.85, rely=0.22, relheight=0.04, relwidth=0.075)

        next_button = tk.Button(self, text="Next slide", command=self.next_ppt)
        next_button.place(relx=0.925, rely=0.22, relheight=0.04, relwidth=0.075)

        speak_button = tk.Button(self, text="Speak", command=self.nao_speak)
        speak_button.place(relx=0.85, rely=0.26, relheight=0.04, relwidth=0.15)

        quit_button = tk.Button(self, text="Quit PPT Mode", command=self.close_connection)
        quit_button.place(relx=0.85, rely=0.30, relheight=0.04, relwidth=0.15)

        chinese_button = tk.Button(self, text="中文", command=self.language_setc)
        chinese_button.place(relx=0.925, rely=0.34, relheight=0.04, relwidth=0.075)

        english_button = tk.Button(self, text="English", command=self.language_sete)
        english_button.place(relx=0.85, rely=0.34, relheight=0.04, relwidth=0.075)

        increase_speech_button = tk.Button(self, text="語速增加", command=self.speech_increase)
        increase_speech_button.place(relx=0.925, rely=0.38, relheight=0.04, relwidth=0.075)

        reduce_speech_button = tk.Button(self, text="語速減少", command=self.speech_reduce)
        reduce_speech_button.place(relx=0.85, rely=0.38, relheight=0.04, relwidth=0.075)

        self.ppt_label = tk.Label(self, textvariable=self.var, bg='white')
        self.ppt_label.place(relheight=1, relwidth=0.85)

        self.text_label = tk.Label(self, textvariable=self.var, bg='gray',
                                   wraplength=150, anchor='nw', justify='left')
        self.text_label.place(relx=0.85, rely=0.42, relheight=0.58, relwidth=0.15)

    def set_img_data(self):
        self.save_dir = tkFileDialog.askdirectory(title="文件儲存的資料夾")
        print(self.save_dir)
        if len(self.save_dir) <= 0:
            tm.showinfo(title='wrong', message='請選擇儲存的資料夾')
            pass
        else:
            if len(self.save_dir) > 0:
                fstart = self.save_dir.rindex('/')
                self.fname = self.save_dir[fstart + 1:]

            self.save_dir = self.save_dir.replace('/', r'\\')
            import glob

            target = self.save_dir + "\\*.jpg"

            self.ppt_len = len(glob.glob(target))
            self.normal_mode = False
            tm.showinfo(title='完成', message='PPT選擇完成')

    def text_set(self, text_list):
        set_text = ''
        for i, text in enumerate(text_list):
            set_text += str(i + 1) + '.\n' + str(text) + '\n'

        self.var.set(set_text)

        if int(self.text_label.winfo_width() / 13) >= 14:
            fontsize = 14
        else:
            fontsize = 10
        self.text_label.config(font=("Helvetica", fontsize, 'bold'), wraplength=self.text_label.winfo_width() - 5)

    def set_data(self):

        self.normal_mode = True
        default_dir = r"C:\Users"
        self.ppt_file = tkFileDialog.askopenfilename(title="選擇文件",
                                                     initialdir=(os.path.expanduser(default_dir)))

        self.save_dir = tkFileDialog.askdirectory(title="儲存的資料夾")

        if len(self.ppt_file) <= 0 or len(self.save_dir) <= 0:
            tm.showinfo(title='wrong', message='請選擇文件及儲存的資料夾')
            pass
        else:
            if len(self.ppt_file) > 0:
                fstart = self.ppt_file.rindex('/')
                self.fname = self.ppt_file[fstart + 1:-5]

            self.ppt_file = self.ppt_file.replace('/', r'\\')
            self.save_dir = self.save_dir.replace('/', r'\\')

            self.ppt2png(self.fname, self.ppt_file, self.save_dir)
            tm.showinfo(title='完成', message='PPT轉換完畢')

    def link_nao(self):

        time.sleep(1)
        ip = self.ip_text.get("1.0", "end-1c")
        print (ip)
        print(self.robot_port)
        if len(ip) <= 0:
            tm.showinfo(title='wrong', message='請輸入ip位址')
        else:
            try:
                self.sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
                self.sock.connect((ip, int(self.robot_port)))
                self.sock.send('Successful connection')
                self.link.set('已連線')
            except TypeError:
                self.link.set('連線失敗')
            except socket:
                self.link.set('連線失敗')

    def img_set(self):
        w_box = self.ppt_label.winfo_width()
        h_box = self.ppt_label.winfo_height()
        if self.normal_mode:
            img_location = self.save_dir + r"\\" + self.fname + r"\\投影片" + str(self.ppt_number + 1) + r".jpg"
            self.text_set(self.ppt_text[self.ppt_number])
        else:
            img_location = self.save_dir + r"\\" + r"\\投影片" + str(self.ppt_number + 1) + r".jpg"
            #self.text_set(self.ppt_text[self.ppt_number])
        print(img_location, type(img_location))
        pil_image = Image.open(img_location)
        pil_image_resized = self.resize(w_box, h_box, pil_image)

        tk_image = ImageTk.PhotoImage(pil_image_resized)
        self.ppt_label.config(image=tk_image)
        self.ppt_label.image = tk_image

    def ppt_start(self):

            self.jpg_dir = self.save_dir + r"\\" + self.fname + r"\\"
            self.ppt_number = 0
            self.img_set()

    def previous_ppt(self):
        if len(self.save_dir) == 0 or len(self.fname) == 0:
            tm.showinfo(title='wrong', message='沒有投影片')
            pass
        elif self.ppt_number == 0:
            tm.showinfo(title='wrong', message='沒有上一頁投影片')
            pass
        else:
            self.ppt_number -= 1
            self.img_set()

    def next_ppt(self):
        if len(self.save_dir) == 0 or len(self.fname) == '':
            tm.showinfo(title='wrong', message='沒有投影片')
            pass
        elif self.ppt_number == self.ppt_len - 1:
            tm.showinfo(title='wrong', message='已經是最後一頁投影片')
            pass
        else:
            self.ppt_number += 1
            self.img_set()

    def ppt2png(self, ppt_file_name, ppt_path, save_path):

        powerpoint = win32com.client.Dispatch('PowerPoint.Application')
        powerpoint.Visible = True
        ppt = powerpoint.Presentations.Open(ppt_path)
        ppt.SaveAs(save_path + r'\\' + ppt_file_name.rsplit('.')[0], 17)
        ppt.Close()
        powerpoint.Quit()

        prs = Presentation(self.ppt_file)
        self.ppt_len = len(prs.slides)

        for slide in prs.slides:
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = ''
                    for run in paragraph.runs:
                        text += run.text
                    text_runs.append(text)
            self.ppt_text.append(text_runs)

    def resize(self, w_box, h_box, pil_image):
        w, h = pil_image.size
        f1 = 1.0 * w_box / w
        f2 = 1.0 * h_box / h
        factor = min([f1, f2])
        width = int(w * factor)
        height = int(h * factor)
        return pil_image.resize((width, height), Image.ANTIALIAS)

    def nao_speak(self):
        if self.link.get() == '關閉連線' or self.link.get() == '連線失敗' or self.link.get() == "尚未連線":
            tm.showinfo(title='wrong', message='請連接機器人後才能說話')
        else:
            texts = self.ppt_text[self.ppt_number]
            send_text = ''
            for text in texts:
                send_text += text + ":::"

            self.sock.send(send_text[:-3])

    def close_connection(self):
        try:
            self.sock.send('Quit ppt mode')
            self.sock.close()
            self.link.set('關閉連線')
        except Exception as e:
            tm.showinfo(title='wrong', message='無法關閉')

    def language_setc(self):
        if self.link.get() == '關閉連線' or self.link.get() == '連線失敗' or self.link.get() == "尚未連線":
            tm.showinfo(title='wrong', message='尚未連線機器人')
        else:
            self.sock.send('Chinese')

    def language_sete(self):
        if self.link.get() == '關閉連線' or self.link.get() == '連線失敗' or self.link.get() == "尚未連線":
            tm.showinfo(title='wrong', message='尚未連線機器人')
        else:
            self.sock.send('English')

    def speech_increase(self):
        if self.link.get() == '關閉連線' or self.link.get() == '連線失敗' or self.link.get() == "尚未連線":
            tm.showinfo(title='wrong', message='尚未連線機器人')
        else:
            self.sock.send('Increased speech rate')

    def speech_reduce(self):
        if self.link.get() == '關閉連線' or self.link.get() == '連線失敗':
            tm.showinfo(title='wrong', message='尚未連線機器人')
        else:
            self.sock.send('Reduce speech rate')


root = tk.Tk()
root.geometry("1280x960")
root.minsize(1024, 768)
window = PPTGUI(root)
root.mainloop()
